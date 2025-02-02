import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
import plotly.io as pio
from pptx import Presentation
from pptx.util import Inches
import os
import base64

# Title of the app
st.title("📊 Data Visualization Dashboard")

# File uploader
uploaded_file = st.file_uploader("Upload a CSV file", type=["csv"])

# Initialize session state for storing charts
if 'charts' not in st.session_state:
    st.session_state.charts = []

if uploaded_file is not None:
    # Load the dataset
    df = pd.read_csv(uploaded_file)
    st.success("File uploaded successfully!")

    # Display the dataset
    st.subheader("Preview of the Dataset")
    st.write(df.head())

    # Sidebar for customization
    st.sidebar.header("Visualization Settings")

    # Select chart type
    chart_type = st.sidebar.selectbox(
        "Select Chart Type",
        ["Bar Chart", "Line Chart", "Scatter Plot", "Histogram", "Box Plot", "Pie Chart"]
    )

    # Select X and Y axes
    x_axis = st.sidebar.selectbox("Select X-axis", df.columns)
    y_axis = st.sidebar.selectbox("Select Y-axis", df.columns)

    # Color picker for plots
    color = st.sidebar.color_picker("Pick a color for the plot", "#1f77b4")

    # Checkbox for showing data points in scatter plot
    show_data_points = st.sidebar.checkbox("Show Data Points in Scatter Plot", value=False)

    # Check if selected axes are valid
    if x_axis and y_axis:
        # Optional: Add a slider for filtering data
        if df[x_axis].dtype in [int, float]:
            min_val, max_val = st.sidebar.slider(
                f"Filter {x_axis} Range",
                float(df[x_axis].min()),
                float(df[x_axis].max()),
                (float(df[x_axis].min()), float(df[x_axis].max()))
            )
            df = df[(df[x_axis] >= min_val) & (df[x_axis] <= max_val)]

        # Button to generate and store the selected chart
        if st.sidebar.button("Generate Chart"):
            st.subheader(f"{chart_type} Visualization")
            try:
                if chart_type == "Bar Chart":
                    fig, ax = plt.subplots()
                    sns.barplot(data=df, x=x_axis, y=y_axis, ax=ax, color=color)
                    st.pyplot(fig)
                elif chart_type == "Line Chart":
                    fig = px.line(df, x=x_axis, y=y_axis, title=f"{y_axis} over {x_axis}", line_shape='linear', color_discrete_sequence=[color])
                    st.plotly_chart(fig)
                elif chart_type == "Scatter Plot":
                    fig = px.scatter(df, x=x_axis, y=y_axis, title=f"{y_axis} vs {x_axis}", color_discrete_sequence=[color])
                    if show_data_points:
                        fig.add_scatter(x=df[x_axis], y=df[y_axis], mode='markers', marker=dict(color=color))
                    st.plotly_chart(fig)
                elif chart_type == "Histogram":
                    fig, ax = plt.subplots()
                    sns.histplot(data=df, x=x_axis, kde=True, ax=ax, color=color)
                    st.pyplot(fig)
                elif chart_type == "Box Plot":
                    fig, ax = plt.subplots()
                    sns.boxplot(data=df, x=x_axis, y=y_axis, ax=ax, color=color)
                    st.pyplot(fig)
                elif chart_type == "Pie Chart":
                    fig = px.pie(df, names=x_axis, values=y_axis, title=f"{y_axis} Distribution by {x_axis}", color_discrete_sequence=[color])
                    st.plotly_chart(fig)

                # Store the chart in session state
                st.session_state.charts.append((chart_type, x_axis, y_axis, color))
                st.success(f"{chart_type} added to the dashboard!")
            except Exception as e:
                st.error(f"Error generating visualization: {e}")

    # Display all generated charts
    if st.session_state.charts:
        st.subheader("Generated Charts")
        for chart in st.session_state.charts:
            chart_type, x_axis, y_axis, color = chart
            st.write(f"{chart_type} - X: {x_axis}, Y: {y_axis}, Color: {color}")

    # Download button for filtered dataset
    st.sidebar.subheader("Download Filtered Dataset")
    csv = df.to_csv(index=False).encode('utf-8')
    st.sidebar.download_button(
        label="Download CSV",
        data=csv,
        file_name='filtered_dataset.csv',
        mime='text/csv'
    )

    # Download dashboard as PowerPoint
    if st.sidebar.button("Download Dashboard as PPT"):
        # Create a PowerPoint presentation
        prs = Presentation()

        for idx, chart in enumerate(st.session_state.charts):
            chart_type, x_axis, y_axis, color = chart

            # Add a slide for each chart
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
            title = slide.shapes.title
            title.text = f"{chart_type} - X: {x_axis}, Y: {y_axis}"

            # Save the chart as an image
            fig, ax = plt.subplots()
            if chart_type == "Bar Chart":
                sns.barplot(data=df, x=x_axis, y=y_axis, ax=ax, color=color)
            elif chart_type == "Line Chart":
                sns.lineplot(data=df, x=x_axis, y=y_axis, ax=ax, color=color)
            elif chart_type == "Scatter Plot":
                sns.scatterplot(data=df, x=x_axis, y=y_axis, ax=ax, color=color)
            elif chart_type == "Histogram":
                sns.histplot(data=df, x=x_axis, kde=True, ax=ax, color=color)
            elif chart_type == "Box Plot":
                sns.boxplot(data=df, x=x_axis, y=y_axis, ax=ax, color=color)
            elif chart_type == "Pie Chart":
                fig = px.pie(df, names=x_axis, values=y_axis, title=f"{y_axis} Distribution by {x_axis}", color_discrete_sequence=[color])
                pio.write_image(fig, f"chart_{idx}.png")
                plt.close(fig)
                continue

            # Save the Matplotlib figure as an image
            plt.savefig(f"chart_{idx}.png")
            plt.close(fig)

            # Add the image to the slide
            left = top = Inches(1)
            slide.shapes.add_picture(f"chart_{idx}.png", left, top, width=Inches(6), height=Inches(4))

        # Save the PowerPoint presentation
        ppt_filename = "dashboard.pptx"
        prs.save(ppt_filename)

        # Provide a download link for the PowerPoint file
        with open(ppt_filename, "rb") as file:
            btn = st.download_button(
                label="Download PPT",
                data=file,
                file_name=ppt_filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

else:
    st.info("Please upload a CSV file to get started.")